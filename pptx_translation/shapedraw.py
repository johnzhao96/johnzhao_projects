"""
Jiayang Zhao
shapedraw: Main module containing functions and classes for drawing pptx shapes.
"""

from pptx.util import Inches, Emu, Pt, Centipoints
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.shapes.shapetree import SlideShapeFactory

from shapetext import draw_text
from shapeutil import emu_to_px, str_emu_to_px, str_cpt_to_px
from shapecolor import extract_color_from_format, get_fill_color, ColorMap

import html

## ************************* Main Shape Drawing Function ************************* ##

def draw_shape(parent_html, pres, slide, shape, color_map, ppi, off_x = 0, off_y = 0, scl_x = 1, scl_y = 1, rot = 0):
    """
    Main shape drawing function. Takes a shape and creates an SVG node representing the
    shape, placing the SVG node as a child of parent_html.
    """
    sbasename = get_basename(shape)
    if sbasename == 'Group':
        draw_group(parent_html, pres, slide, shape, color_map, ppi, off_x, off_y, scl_x, scl_y, rot)
    else:
        sleft = shape.left + off_x
        stop = shape.top + off_y
        swidth = shape.width * scl_x
        sheight = shape.height * scl_y
        
        curr_html = parent_html

        if sbasename == 'Straight Connector':
            curr_html.line('', 
                x1 = str_emu_to_px(sleft, ppi), 
                y1 = str_emu_to_px(stop, ppi), 
                x2 = str_emu_to_px(sleft + swidth, ppi), 
                y2 = str_emu_to_px(stop + sheight, ppi),
                transform = get_transform(shape, ppi, off_x, off_y, scl_x, scl_y),
                style = get_style(shape, color_map, ppi)
            )
        elif sbasename == 'Oval':
            curr_html.ellipse('', 
                cx = str_emu_to_px(sleft + swidth/2, ppi), 
                cy = str_emu_to_px(stop + sheight/2, ppi), 
                rx = str_emu_to_px(swidth/2, ppi), 
                ry = str_emu_to_px(sheight/2, ppi), 
                transform = get_transform(shape, ppi, off_x, off_y, scl_x, scl_y),
                style = get_style(shape, color_map, ppi)
            )
        elif sbasename == 'Rectangle' or sbasename == 'Rounded Rectangle':
            curr_html.rect('', 
                x = str_emu_to_px(sleft, ppi), 
                y = str_emu_to_px(stop, ppi), 
                width = str_emu_to_px(swidth, ppi), 
                height = str_emu_to_px(sheight, ppi), 
                transform = get_transform(shape, ppi, off_x, off_y, scl_x, scl_y),
                style = get_style(shape, color_map, ppi)
            )
        else:
            curr_html.rect('', 
                x = str_emu_to_px(sleft, ppi), 
                y = str_emu_to_px(stop, ppi), 
                width = str_emu_to_px(swidth, ppi), 
                height = str_emu_to_px(sheight, ppi), 
                transform = get_transform(shape, ppi, off_x, off_y, scl_x, scl_y),
                style = get_style(shape, color_map, ppi)
            )
        
        if shape.has_text_frame and shape.text_frame.text != '':
            text_box = curr_html.foreignObject(
                x = str_emu_to_px(sleft, ppi),
                y = str_emu_to_px(stop, ppi),
                width = str_emu_to_px(swidth, ppi),
                height = str_emu_to_px(sheight, ppi),
                transform = get_transform(shape, ppi, off_x, off_y, scl_x, scl_y)
            )
            draw_text(text_box, pres, slide, shape, color_map, ppi)

def draw_group(parent_html, pres, slide, shape, color_map, ppi, off_x, off_y, scl_x, scl_y, rot):
    group = GroupShape(shape)

    # Find bounding box location of shapes
    if len(group.children) == 0:
        return

    if group.chdcx != 0:
        scale_x = group.cx / group.chdcx
    else:
        scale_x = 1

    if group.chdcy != 0:
        scale_y = group.cy / group.chdcy
    else:
        scale_y = 1
    
    offset_x = group.x - group.chdx
    offset_y = group.y - group.chdy

    for shp in group.children:
        shp_offset_x = offset_x + off_x + (shp.left - group.chdx) * (scale_x - 1) 
        shp_offset_y = offset_y + off_y + (shp.top - group.chdy) * (scale_y - 1)
        shp_scale_x = scale_x*scl_x
        shp_scale_y = scale_y*scl_y

        draw_shape(parent_html, pres, slide, shp, color_map, ppi, shp_offset_x, shp_offset_y, shp_scale_x, shp_scale_y, rot)

class GroupShape:
    """
    Class for dealing with Group Shapes
    """
    def __init__(self, shape):
        self.shape = shape
        self.children = []
        try:
            for child in shape.element.iter_shape_elms():
                chd = SlideShapeFactory(child, shape)
                self.children.append(chd)
        except AttributeError:
            print("Attribute Error: {}".format(shape.name))
            return None
        
        self.x = shape.left
        self.y = shape.top
        self.cx = shape.width
        self.cy = shape.height

        chdx_lst = shape.element.xpath('./p:grpSpPr/a:xfrm/a:chOff/@x')
        if len(chdx_lst) > 0:
            self.chdx = Emu(float(chdx_lst[0]))

        chdy_lst = shape.element.xpath('./p:grpSpPr/a:xfrm/a:chOff/@y')
        if len(chdy_lst) > 0:
            self.chdy = Emu(float(chdy_lst[0]))

        chdcx_lst = shape.element.xpath('./p:grpSpPr/a:xfrm/a:chExt/@cx')
        if len(chdcx_lst) > 0:
            self.chdcx = Emu(float(chdcx_lst[0]))

        chdcy_lst = shape.element.xpath('./p:grpSpPr/a:xfrm/a:chExt/@cy')
        if len(chdcy_lst) > 0:
            self.chdcy = Emu(float(chdcy_lst[0]))

        self.rot = shape.rotation


## ************************* Shape Property Getters ************************* ##

def get_style(shape, color_map, ppi):
    """
    Given a shape, gets the shape's border (line) style and fill formatting
    """
    return '{} {}'.format(get_border_style(shape, color_map, ppi), get_fill_color(shape, color_map))

def get_border_style(shape, color_map, ppi):
    """
    Given a shape, returns a string with the border/fill style of the shape
    """
    try:
        shape_line = shape.line    
    except AttributeError:
        return 'stroke:black ; stroke-width:1 ;'
    stroke_width = str_emu_to_px(shape_line.width, ppi)
    stroke = extract_color_from_format(shape_line.fill, color_map)
    if stroke == 'no color':
        stroke_width = 0
    return 'stroke:{} ; stroke-width:{} ;'.format(stroke, stroke_width)

def get_basename(shape):
    """
    Given a shape name, retrieves the shape's base name
    """
    name = shape.name
    index = -1
    n = -1 * len(name)
    while name[index] != ' ':
        index -= 1
        if index <= n:
            return name
    return name[0:index]


## ************************* Shape Transformations ************************* ##

def get_flipV(shape):
    """
    Gets the flipV (vertical flip) attribute of a shape (if it has one). Returns True if flipV = '1' 
    and False otherwise
    """
    flip_list = shape.element.xpath('./p:spPr/a:xfrm/@flipV')
    if len(flip_list) > 0:
        if flip_list[0] == '1':
            return True
    return False

def get_flipH(shape):
    """
    Gets the flipH (horizontal flip) attribute of a shape (if it has one). Returns True if flipH = '1' 
    and False otherwise
    """
    flip_list = shape.element.xpath('./p:spPr/a:xfrm/@flipH')
    if len(flip_list) > 0:
        if flip_list[0] == '1':
            return True
    return False

def get_rotation(shape):
    """
    Gets the rotation value of a shape in degrees
    """
    return shape.rotation


def get_transform(shape, ppi, off_x, off_y, scl_x, scl_y):
    """
    Gets any transformation style attributes (flips, rotations, etc) that should be applied to the shape SVG
    """
    sflipH = get_flipH(shape)
    sflipV = get_flipV(shape)
    srot = get_rotation(shape)
    
    if sflipH or sflipV or srot != 0:
        if sflipH:
            fH = -1
        else:
            fH = 1
        
        if sflipV:
            fV = -1
        else:
            fV = 1

        sx = shape.left + off_x
        sy = shape.top + off_y
        scx = shape.width * scl_x
        scy = shape.height * scl_y

        tx = emu_to_px(sx + (scx / 2), ppi)
        ty = emu_to_px(sy + (scy / 2), ppi)

        # Changes coordinate system such that the origin is in the middle of the shape,
        # performs the transformation, then changes coordinates back to the original system
        return 'translate({} {}) scale({} {}) rotate({}) translate({} {})'.format(tx, ty, fH, fV, srot, -1 * tx, -1 * ty)
    return ''
