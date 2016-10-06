#!/usr/bin/env python

"""
Jiayang Zhao
pptx_to_html: Script for parsing a PowerPoint presentation and rendering the presentation in HTML.
"""

from pptx import Presentation
from pptx.util import Inches, Emu
from html import HTML

import shapedraw
from shapecolor import construct_pres_color_map

import sys
from os.path import basename, join

# Pixels Per Inch
PPI = 50
slideshow_PPI = 75

slideshow_mode = (sys.argv[1] == '--slideshow')
    
if slideshow_mode:
    offset = 1
else:
    offset = 0

inpath = sys.argv[1+offset]
outpath = sys.argv[2+offset]
prs = Presentation(inpath)

SLIDE_HEIGHT = prs.slide_height
SLIDE_WIDTH = prs.slide_width 
COLOR_MAP = construct_pres_color_map(prs)

def draw_slide(parent_html, pres, slide, display_w, display_h, ppi):
    """
    Outputs an entire slide as an HTML with children nodes representing each of the slide's shapes
    """
    slide_svg = parent_html.svg('',
        xmlns='http://www.w3.org/2000/svg', 
        width='{}px'.format(str(display_w)), 
        height='{}px'.format(str(display_h)), 
        x='0', 
        y='0', 
        viewbox='0 0 {} {}'.format(str(display_w), str(display_h)),
        newlines=True
    )
    slide_svg.rect('',
        x='0',
        y='0',
        width=str(display_w),
        height=str(display_h),
        style='fill:white ; stroke:black ; stroke-width:1'
    )
    for shape in slide.shapes:
        shapedraw.draw_shape(slide_svg, pres, slide, shape, COLOR_MAP, ppi)    
    slide_svg.rect('',
        x='0',
        y='0',
        width=str(display_w),
        height=str(display_h),
        style='fill-opacity:0 ; stroke:black ; stroke-width:1'
    )

if slideshow_mode:
    dwidth = Emu(SLIDE_WIDTH).inches * slideshow_PPI
    dheight = Emu(SLIDE_HEIGHT).inches * slideshow_PPI

    infilebase = basename(inpath)[:-5]

    for n,slide in enumerate(prs.slides):
        my_html = HTML('html', newlines=True)
        my_html.a('Previous Slide', href='log_{}_{}.html'.format(infilebase, n))
        draw_slide(my_html, prs, slide, dwidth, dheight, slideshow_PPI)
        my_html.a('Next Slide', href='log_{}_{}.html'.format(infilebase, n+2))

        destpath = join(outpath, 'log_{}_{}.html'.format(infilebase, n+1))
        outfile = open(destpath, 'w')
        outfile.write(str(my_html))
        outfile.close()

else: 
    my_html = HTML('html', newlines=True)

    dwidth = Emu(SLIDE_WIDTH).inches * PPI
    dheight = Emu(SLIDE_HEIGHT).inches * PPI

    for slide in prs.slides:
        draw_slide(my_html, prs, slide, dwidth, dheight, PPI)
        
    outfile = open(outpath, 'w')
    outfile.write(str(my_html))
    outfile.close()
    
