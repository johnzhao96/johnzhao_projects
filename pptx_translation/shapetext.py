"""
Jiayang Zhao
shapetext: Module for parsing shape text and formattings.
"""

from pptx.util import Inches, Emu, Pt, Centipoints
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

from shapeutil import emu_to_px, str_emu_to_px, str_cpt_to_px, replace_spaces
from shapecolor import extract_color_from_format, get_fill_color


## ************************* Main Text Drawing Function ************************* ##

def draw_text(text_html, pres, slide, shape, color_map, ppi):
    """
    Main text drawing function. Given a shape, outputs its text as a textbox that lies on top
    of the shape SVG.
    """
    text_html.meta(charset = 'utf-8')
    text_box = text_html.table('', 
        style = 'border-collapse: collapse; width: 100%; height: 100%;',
        newlines = True
    ).tr.td('',
        style = '{}'.format(get_vertical_alignment(pres, slide, shape))
    )
    for paragraph in shape.text_frame.paragraphs:
        para_html = text_box.p('', 
            xmlns = 'http://www.w3.org/1999/xhtml', 
            style = '{} ; {}'.format(
                get_margins(shape.text_frame, ppi), 
                get_alignment(pres, slide, shape, paragraph)
            ),
            newlines = False
        )
        for run in paragraph.runs:
            font = Font(pres, slide, shape, paragraph, run, color_map)
            stext = replace_spaces(run.text).encode('utf-8')
            curr_node = para_html
            if font.bold == True:
                curr_node = curr_node.b
            if font.italic == True:
                curr_node = curr_node.i
            if font.underline == True:
                curr_node = curr_node.u
            if font.size is not None:
                fontsize = 'font-size:{}'.format(str_emu_to_px(font.size, ppi))
            else:
                fontsize = ''
            if font.name is not None:
                fontname = 'font-family:{}'.format(font.name.encode('utf-8'))
            else:
                fontname = ''
            curr_node.span(stext, style='{} ; {} ; color:{}'.format(fontsize, fontname, font.color), escape = False)
            sibling_list = run._r.xpath('./following-sibling::*')
            if len(sibling_list) > 0:
                for sib in sibling_list:
                    if sib.tag == '{http://schemas.openxmlformats.org/drawingml/2006/main}br':
                        curr_node.br
                    else:
                        break

class Font:
    """
    Class for dealing with text Font formatting
    """
    def __init__(self, pres, slide, shape, paragraph, run, color_map):
        self.size = init_font_size(pres, slide, shape, paragraph, run)
        self.name = init_font_name(pres, slide, shape, paragraph, run)

        self.bold = run.font.bold
        self.italic = run.font.italic
        self.underline = run.font.underline
        self.color = extract_color_from_format(run.font.fill, color_map)
        
        if self.bold is None:
            self.bold = paragraph.font.bold
        if self.italic is None:
            self.italic = paragraph.font.italic
        if self.underline is None:
            self.underline = paragraph.font.underline
        if self.color == 'no color':
            self.color = extract_color_from_format(paragraph.font.fill, color_map)
            if self.color == 'no color':
                self.color = 'black'


## ************************* Placeholder Getters ************************* ##

def find_layout_placeholder(slide, shape):
    """
    Returns the slide layout placeholder shape to which a shape corresponds (if the shape is a placeholder shape)
    """
    sidx = shape.placeholder_format.idx
    try:
        return slide.slide_layout.placeholders[sidx]
    except IndexError: 
        ph_type = shape.placeholder_format.type
        for ph in slide.slide_layout.placeholders:
            if ph.placeholder_format.type == ph_type:
                return ph
        return None
    
def find_master_placeholder(slide, shape):
    """
    Returns the slide master placeholder shape to which a shape corresponds (if the shape is a placeholder shape)
    """
    sidx = shape.placeholder_format.idx
    try:
        return slide.slide_layout.slide_master.placeholders[sidx]
    except IndexError: 
        ph_type = shape.placeholder_format.type
        for ph in slide.slide_layout.slide_master.placeholders:
            if ph.placeholder_format.type == ph_type:
                return ph
        return None


## ************************* Get Font Size ************************* ##

def init_font_size(pres, slide, shape, paragraph, run):
    """
    Gets a the font size of one run in the text of a shape
    """
    self_size = run.font.size
    lvl = paragraph.level

    # inherit from paragraph if not defined at run level
    if self_size is None:
        self_size = paragraph.font.size

    if self_size is None:
        self_size = get_font_size_from_lstStyle(shape, lvl)

    # inherit from placeholder if needed
    if shape.is_placeholder:
        if self_size is None:
            sl_ph = find_layout_placeholder(slide, shape)
            if not sl_ph is None:
                self_size = get_font_size_from_lstStyle(sl_ph, lvl) 
        if self_size is None:
            sm_ph = find_master_placeholder(slide, shape)
            if not sm_ph is None:
                self_size = get_font_size_from_lstStyle(sm_ph, lvl)
        if self_size is None:
            self_size = get_font_size_from_master(slide, shape, lvl)

    # inherit from presentation defaults
    if self_size is None:
        self_size = get_default_font_size(pres, lvl)

    return self_size

def get_font_size_from_lstStyle(ph_shape, lvl):
    """
    Returns font size from placeholder style. Font sizes returned in EMU
    """
    lvl_elm_name = 'a:lvl{}pPr'.format(lvl+1)
    sz_list = ph_shape.element.xpath('./p:txBody/a:lstStyle/{}/a:defRPr/@sz'.format(lvl_elm_name))
    if len(sz_list) > 0:
        return Centipoints(float(sz_list[0])).emu
    else:
        return None

def get_font_size_from_master(slide, ph_shape, lvl):
    """
    Gets the slide master default font size for a particular shape's text
    """
    lvl_elm_name = 'a:lvl{}pPr'.format(lvl+1)
    if ph_shape.placeholder_format.type in [PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.VERTICAL_TITLE]:
        style_type = 'p:titleStyle'
    elif ph_shape.placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.VERTICAL_BODY]:
        style_type = 'p:bodyStyle'
    else:
        style_type = 'p:otherStyle'
    smaster = slide.slide_layout.slide_master
    sz_list = smaster._element.xpath('./p:txStyles/{}/{}/a:defRPr/@sz'.format(style_type, lvl_elm_name))
    if len(sz_list) > 0:
        return Centipoints(float(sz_list[0])).emu
    else:
        return None

def get_default_font_size(pres, lvl):
    """
    Gets the default font size of the entire PPTX document
    """
    """
    pres_xml = pres._presentation._element
    lvl_elm_name = 'a:lvl{}pPr'.format(lvl+1)
    sz_list = pres_xml.xpath('./p:defaultTextStyle/{}/a:defRPr/@sz'.format(lvl_elm_name))
    if len(sz_list) > 0:
        return Centipoints(float(sz_list[0])).emu
    else:
        return None
    """
    return None


## ************************* Get Horizontal Alignment ************************* ##

# Gets text alignment information
def get_alignment(pres, slide, shape, paragraph):
    """
    Gets a shape's text's paragraph's horizontal alignemtn
    """
    ret = ''
    lvl = paragraph.level

    alignment = paragraph.alignment

    if alignment is None:
        alignment = get_font_align_from_lstStyle(shape, lvl)

    if shape.is_placeholder:
        if alignment is None:
            sl_ph = find_layout_placeholder(slide, shape)
            if not sl_ph is None:
                alignment = get_font_align_from_lstStyle(sl_ph, lvl) 
        if alignment is None:
            sm_ph = find_master_placeholder(slide, shape)
            if not sm_ph is None:
                alignment = get_font_align_from_lstStyle(sm_ph, lvl)
        if alignment is None:
            alignment = get_font_align_from_master(slide, shape, lvl)

    if alignment == PP_ALIGN.CENTER:
        ret = 'text-align:center'
    elif alignment == PP_ALIGN.JUSTIFY or alignment == PP_ALIGN.JUSTIFY_LOW:
        ret = 'text-align:justify'
    elif alignment == PP_ALIGN.LEFT:
        ret = 'text-align:left'
    elif alignment == PP_ALIGN.RIGHT:
        ret = 'text-align:right'
    return ret

def parse_align_attr(algn_str):
    """
    Parses a horizontal alignment attribute from the PPTX XML and returns the corresponding PP_ALIGN enumeration
    """
    if algn_str in ALIGNMENT_ENUMS:
    	return ALIGNMENT_ENUMS[algn_str]
    return None

ALIGNMENT_ENUMS = {
	'ctr' : PP_ALIGN.CENTER,
	'just' : PP_ALIGN.JUSTIFY,
	'justlow' : PP_ALIGN.JUSTIFY_LOW,
	'l' : PP_ALIGN.LEFT,
	'r' : PP_ALIGN.RIGHT,
	'dist' : PP_ALIGN.DISTRIBUTE,
	'thaiDist' : PP_ALIGN.THAI_DISTRIBUTE
}

def get_font_align_from_lstStyle(ph_shape, lvl):
    """
    Parses a text box's placeholder and returns the placeholder default horizontal alignment
    """
    lvl_elm_name = 'a:lvl{}pPr'.format(lvl+1)
    algn_list = ph_shape.element.xpath('./p:txBody/a:lstStyle/{}/@algn'.format(lvl_elm_name))
    if len(algn_list) > 0:
        algn_str = algn_list[0]
        return parse_align_attr(algn_str)
    else:
        return None

def get_font_align_from_master(slide, ph_shape, lvl):
    """
    Parses a shape's presentation's slide master and returns the master default horizontal alignment
    """
    lvl_elm_name = 'a:lvl{}pPr'.format(lvl+1)
    if ph_shape.placeholder_format.type in [PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.VERTICAL_TITLE]:
        style_type = 'p:titleStyle'
    elif ph_shape.placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.VERTICAL_BODY]:
        style_type = 'p:bodyStyle'
    else:
        style_type = 'p:otherStyle'
    smaster = slide.slide_layout.slide_master
    algn_list = smaster._element.xpath('./p:txStyles/{}/{}/@algn'.format(style_type, lvl_elm_name))
    if len(algn_list) > 0:
        algn_str = algn_list[0]
        return parse_align_attr(algn_str)
    else:
        return None


## ************************* Get Vertical Alignment ************************* ##

def get_vertical_alignment(pres, slide, shape):
    """
    Gets a shape's text vertical alignment
    """
    ret = ''

    valignment = shape.text_frame.vertical_anchor

    if shape.is_placeholder:
        if valignment is None:
            sl_ph = find_layout_placeholder(slide, shape)
            if not sl_ph is None:
                valignment = get_font_valign_from_txBody(sl_ph) 
        if valignment is None:
            sm_ph = find_master_placeholder(slide, shape)
            if not sm_ph is None:
                valignment = get_font_valign_from_txBody(sm_ph)

    if valignment == MSO_VERTICAL_ANCHOR.TOP:
        ret = 'vertical-align:top'
    elif valignment == MSO_VERTICAL_ANCHOR.BOTTOM:
        ret = 'vertical-align:bottom'
    elif valignment == MSO_VERTICAL_ANCHOR.MIDDLE:
        ret = 'vertical-align:middle'
    return ret

def parse_valign_attr(valgn_str):
    """
    Parses a vertical alignment attribute from the PPTX XML and returns the corresponding MSO_VERTICAL_ANCHOR enumeration
    """
    if valgn_str in VERTICAL_ALIGNMENT_ENUMS:
    	return VERTICAL_ALIGNMENT_ENUMS[valgn_str]
    return None

VERTICAL_ALIGNMENT_ENUMS = {
	't' : MSO_VERTICAL_ANCHOR.TOP,
	'ctr' : MSO_VERTICAL_ANCHOR.MIDDLE,
	'b' : MSO_VERTICAL_ANCHOR.BOTTOM
}

def get_font_valign_from_txBody(ph_shape):
    """
    Parses the anchor (vertical alignment) attribute of a shape's p:txBody node
    """
    valgn_list = ph_shape.element.xpath('./p:txBody/@anchor')
    if len(valgn_list) > 0:
        valgn_str = valgn_list[0]
        return parse_valign_attr(valgn_str)
    else:
        return None


## ************************* Get Font Name ************************* ##

# *** NEED TO DO: inherit font name from upper levels
def init_font_name(pres, slide, shape, paragraph, run):
    """
    Gets a shape's font name
    """
    self_name = run.font.name
    lvl = paragraph.level

    # inherit from paragraph if not defined at run level
    if self_name is None:
        self_name = paragraph.font.name

    return self_name


## ************************* Get Margins ************************* ##

def get_margins(text_frame, ppi):
    """
    Gets text margin information
    """
    margins = 'margin: {} {} {} {}'.format(
        str_emu_to_px(text_frame.margin_top, ppi),
        str_emu_to_px(text_frame.margin_right, ppi),
        str_emu_to_px(text_frame.margin_bottom, ppi),
        str_emu_to_px(text_frame.margin_left, ppi)
    )
    return margins
